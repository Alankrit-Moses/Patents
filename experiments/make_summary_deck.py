from __future__ import annotations

import csv
import json
import re
import textwrap
from collections import defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "experiments" / "artifacts" / "results"
OUT_DIR = ROOT / "experiments" / "artifacts" / "presentation"
OUT_PPTX = OUT_DIR / "framework_experiment_summary_large_text.pptx"
OUT_AUDIT = OUT_DIR / "judge_score_audit.csv"

MODELS = [
    ("gpt-4.1", "gpt-4.1"),
    ("gpt-oss-high", "gpt-oss-high"),
    ("gpt-oss-medium", "gpt-oss-medium"),
    ("qwen3-32B-nonreasoning", "qwen3-32B non-reasoning"),
    ("llama-3.3-70b", "llama3-70B"),
]
SETUPS = ["A", "B-query", "B-resolved", "C"]
SUMMARY_SETUP_KEY = {"A": "A", "B-query": "B", "B-resolved": "B", "C": "C"}
METRICS = [
    ("E1", "experiment_1", "E_text_alignment_scores"),
    ("E2", "experiment_2", "E_text_alignment_score"),
    ("E3-M", "experiment_3", "M_alignment_score"),
    ("E3-T", "experiment_3", "T_alignment_score"),
]

BG = RGBColor(247, 248, 244)
INK = RGBColor(35, 38, 36)
MUTED = RGBColor(92, 99, 94)
ACCENT = RGBColor(42, 111, 108)
ACCENT_DARK = RGBColor(24, 78, 75)
WIN_FILL = RGBColor(222, 239, 177)
HEADER_FILL = RGBColor(37, 72, 76)
PANEL = RGBColor(255, 255, 252)
RULE = RGBColor(206, 211, 202)


def load_jsonl(path: Path) -> list[dict[str, Any]]:
    return [json.loads(line) for line in path.read_text(encoding="utf-8").splitlines() if line.strip()]


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace").strip()


def clean_text(value: str, limit: int = 320) -> str:
    value = re.sub(r"\s+", " ", str(value)).strip()
    if len(value) <= limit:
        return value
    return value[: limit - 1].rstrip() + "..."


def saved_excerpt(path: Path) -> str:
    text = read_text(path)
    parts = re.split(r"\r?\n\s*\r?\n", text, maxsplit=1)
    return parts[1].strip() if len(parts) == 2 else text


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.55))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(29)
    p.font.bold = True
    p.font.color.rgb = INK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.47), Inches(0.78), Inches(12.0), Inches(0.3))
        sf = sub.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Aptos"
        sp.font.size = Pt(15)
        sp.font.color.rgb = MUTED


def add_band(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    title_size: int = 16,
    body_size: int = 19,
) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT_DARK
    for line in body.split("\n"):
        bp = tf.add_paragraph()
        bp.text = line
        bp.font.name = "Aptos"
        bp.font.size = Pt(body_size)
        bp.font.color.rgb = INK
        bp.space_before = Pt(2)


def add_bullets(slide, x: float, y: float, w: float, h: float, items: list[str], size: int = 23) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(7)


def summary_value(model: str, setup: str, pattern: str, experiment: str, metric: str) -> float | None:
    path = RESULTS / model / "robustness" / f"{setup}.summary.json"
    if not path.exists():
        return None
    data = json.loads(path.read_text(encoding="utf-8"))
    key = f"{experiment}/setup_{SUMMARY_SETUP_KEY[setup]}/pattern_{pattern}/{metric}"
    item = data.get("groups", {}).get(key)
    return None if item is None else float(item["mean"])


def model_table_values(model: str, pattern: str | None) -> dict[str, list[float | None]]:
    values: dict[str, list[float | None]] = {}
    for setup in SETUPS:
        row = []
        for _, experiment, metric in METRICS:
            if pattern is None:
                p1 = summary_value(model, setup, "P1", experiment, metric)
                p2 = summary_value(model, setup, "P2", experiment, metric)
                row.append(None if p1 is None or p2 is None else (p1 + p2) / 2)
            else:
                row.append(summary_value(model, setup, pattern, experiment, metric))
        values[setup] = row
    return values


def add_results_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    values: dict[str, list[float | None]],
    table_font_size: float = 20,
    header_font_size: float = 18,
    label_size: float = 16,
) -> None:
    title = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.22))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Aptos"
    p.font.size = Pt(label_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT_DARK

    table = slide.shapes.add_table(5, 5, Inches(x), Inches(y + 0.25), Inches(w), Inches(h - 0.25)).table
    headers = ["Setup", "E1", "E2", "E3-M", "E3-T"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        para = cell.text_frame.paragraphs[0]
        para.font.name = "Aptos"
        para.font.size = Pt(header_font_size)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER

    winners = []
    for idx in range(4):
        col = [row[idx] for row in values.values() if row[idx] is not None]
        winners.append(max(col) if col else None)

    for row_idx, setup in enumerate(SETUPS, start=1):
        row_values = values[setup]
        cells = [setup] + ["n/a" if v is None else f"{v:.2f}" for v in row_values]
        for col_idx, text in enumerate(cells):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            is_winner = col_idx > 0 and row_values[col_idx - 1] is not None and row_values[col_idx - 1] == winners[col_idx - 1]
            cell.fill.fore_color.rgb = WIN_FILL if is_winner else PANEL
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Aptos"
            para.font.size = Pt(table_font_size)
            para.font.bold = bool(is_winner or col_idx == 0)
            para.font.color.rgb = INK
            para.alignment = PP_ALIGN.CENTER if col_idx else PP_ALIGN.LEFT
            cell.margin_left = Inches(0.02)
            cell.margin_right = Inches(0.02)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)


def add_results_slide(prs: Presentation, title: str, pattern: str | None) -> None:
    groups = [MODELS[:3], MODELS[3:]]
    for group_index, group in enumerate(groups, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        subtitle = "Robustness means; highlighted cells are setup winners within each model and experiment."
        add_title(slide, f"{title} ({group_index}/2)", subtitle)
        y = 1.05 if len(group) == 3 else 1.35
        table_h = 1.85 if len(group) == 3 else 2.25
        step = 2.05 if len(group) == 3 else 2.55
        for model, label in group:
            add_results_table(
                slide,
                0.65,
                y,
                12.05,
                table_h,
                label,
                model_table_values(model, pattern),
                table_font_size=20,
                header_font_size=18,
                label_size=16,
            )
            y += step


def add_results_slide_legacy(prs: Presentation, title: str, pattern: str | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = "Robustness means; highlighted cells are setup winners within each model and experiment."
    add_title(slide, title, subtitle)
    y = 1.05
    for model, label in MODELS:
        add_results_table(
            slide,
            0.65,
            y,
            12.05,
            1.12,
            label,
            model_table_values(model, pattern),
            table_font_size=20,
            header_font_size=18,
            label_size=16,
        )
        y += 1.2


def make_result_slides(prs: Presentation) -> None:
    add_results_slide(prs, "Pattern 1 Results", "P1")
    add_results_slide(prs, "Pattern 2 Results", "P2")
    add_results_slide(prs, "Macro Average Results", None)


def task_records() -> dict[str, dict[str, Any]]:
    return {row["task_id"]: row for row in load_jsonl(ROOT / "experiments" / "artifacts" / "tasks.jsonl")}


def result_records(model: str, setup: str) -> dict[str, dict[str, Any]]:
    path = RESULTS / model / "robustness" / f"{setup}.results.jsonl"
    records = {}
    for row in load_jsonl(path):
        key = row.get("sample_id") or f"{row.get('task_id')}:{row.get('sample_index')}"
        records[key] = row
    return records


def audit_rows() -> list[dict[str, Any]]:
    tasks = task_records()
    rows: list[dict[str, Any]] = []
    for model, _ in MODELS:
        for setup in SETUPS:
            eval_path = RESULTS / model / "robustness" / f"{setup}.eval.jsonl"
            if not eval_path.exists():
                continue
            results = result_records(model, setup)
            for evaluation in load_jsonl(eval_path):
                if evaluation.get("errors"):
                    continue
                key = evaluation.get("sample_id") or f"{evaluation.get('task_id')}:{evaluation.get('sample_index')}"
                result = results.get(key)
                task = tasks.get(str(evaluation.get("task_id")))
                if result is None or task is None:
                    continue
                parsed = result.get("parsed_output", {})
                experiment = str(evaluation["experiment"])
                if experiment == "1":
                    examples = parsed.get("examples") if isinstance(parsed, dict) else []
                    examples = examples if isinstance(examples, list) else []
                    pattern_text = read_text(ROOT / task["m_path"])
                    scores = evaluation.get("metrics", {}).get("E_text_alignment_scores", [])
                    for idx, score in enumerate(scores):
                        pred = ""
                        if idx < len(examples) and isinstance(examples[idx], dict):
                            pred = examples[idx].get("excerpt", "")
                        rows.append(
                            {
                                "score": int(score),
                                "model": model,
                                "setup": setup,
                                "experiment": "E1",
                                "pattern": evaluation.get("pattern_id", ""),
                                "task_id": evaluation.get("task_id", ""),
                                "prediction": clean_text(pred, 900),
                                "ground_truth": clean_text("Canonical M: " + pattern_text, 900),
                            }
                        )
                elif experiment == "2":
                    pred = parsed.get("excerpt", "") if isinstance(parsed, dict) else ""
                    gold = saved_excerpt(ROOT / task["gold_e_text_path"])
                    score = evaluation.get("metrics", {}).get("E_text_alignment_score")
                    rows.append(
                        {
                            "score": int(score),
                            "model": model,
                            "setup": setup,
                            "experiment": "E2",
                            "pattern": evaluation.get("pattern_id", ""),
                            "task_id": evaluation.get("task_id", ""),
                            "prediction": clean_text(pred, 900),
                            "ground_truth": clean_text("Reference E.text: " + gold, 900),
                        }
                    )
                elif experiment == "3":
                    for metric, field, label, path_key in [
                        ("M_alignment_score", "M", "E3-M", "m_path"),
                        ("T_alignment_score", "T", "E3-T", "t_path"),
                    ]:
                        score = evaluation.get("metrics", {}).get(metric)
                        if score is None:
                            continue
                        pred = parsed.get(field, "") if isinstance(parsed, dict) else ""
                        gold = read_text(ROOT / task[path_key])
                        rows.append(
                            {
                                "score": int(score),
                                "model": model,
                                "setup": setup,
                                "experiment": label,
                                "pattern": evaluation.get("pattern_id", ""),
                                "task_id": evaluation.get("task_id", ""),
                                "prediction": clean_text(pred, 900),
                                "ground_truth": clean_text("Canonical " + field + ": " + gold, 900),
                            }
                        )
    return rows


def write_audit_csv(rows: list[dict[str, Any]]) -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    with OUT_AUDIT.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=["score", "model", "setup", "experiment", "pattern", "task_id", "prediction", "ground_truth"],
        )
        writer.writeheader()
        writer.writerows(rows)


def choose_audit_examples(
    rows: list[dict[str, Any]], model_filter: str = "gpt-4.1"
) -> dict[int, list[dict[str, Any]]]:
    rows = [row for row in rows if row["model"] == model_filter]
    by_score: dict[int, list[dict[str, Any]]] = defaultdict(list)
    for row in rows:
        by_score[int(row["score"])].append(row)

    chosen: dict[int, list[dict[str, Any]]] = {}
    preferred_experiments = ["E1", "E2", "E3-M", "E3-T"]
    for score in range(6):
        picked: list[dict[str, Any]] = []
        pool = by_score.get(score, [])
        for experiment in preferred_experiments:
            candidates = [row for row in pool if row["experiment"] == experiment and row not in picked]
            if candidates:
                candidates.sort(key=lambda row: len(row["prediction"]) + len(row["ground_truth"]))
                picked.append(candidates[0])
            if len(picked) == 2:
                break
        if len(picked) < 2:
            rest = [row for row in pool if row not in picked]
            rest.sort(key=lambda row: len(row["prediction"]) + len(row["ground_truth"]))
            picked.extend(rest[: 2 - len(picked)])
        chosen[score] = picked
    return chosen


def add_audit_slide(prs: Presentation, score: int, examples: list[dict[str, Any]], count: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    labels = {
        0: "No recoverable grounded excerpt",
        1: "Unrelated or not aligned",
        2: "Weakly related, central conditions missing",
        3: "Partial or uncertain alignment",
        4: "Valid with minor omission or ambiguity",
        5: "Clear and complete alignment",
    }
    add_title(slide, f"Judge Score {score}", f"{labels[score]} | {count} gpt-4.1 audit rows in retained robustness evaluations")
    if not examples:
        add_bullets(slide, 0.7, 1.4, 11.8, 1.0, ["No retained examples found for this score."], 23)
        return
    y = 1.25
    for idx, row in enumerate(examples, start=1):
        h = 2.35
        add_band(
            slide,
            0.55,
            y,
            12.2,
            h,
            f"{idx}. {row['experiment']} | {row['pattern']} | {row['model']} | {row['setup']} | score {row['score']}",
            "Prediction: "
            + clean_text(row["prediction"], 190)
            + "\nGround truth: "
            + clean_text(row["ground_truth"], 190),
            title_size=18,
            body_size=20,
        )
        y += h + 0.3


def add_setup_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Three Setups", "The comparison isolates direct prompting, framework-grounded prompting, and framework-inspired agency.")
    add_band(
        slide,
        0.6,
        1.2,
        3.8,
        4.8,
        "Setup A: baseline",
        "Minimal prose prompt.\nOne model call.\nNo framework notation or decomposition beyond the task statement.",
    )
    add_band(
        slide,
        4.75,
        1.2,
        3.8,
        4.8,
        "Setup B: structured prompting",
        "Framework-grounded single prompt.\nQuery-only: model receives the formal query and must infer the task from the framework rules.\nResolved-task: adds a natural-language resolution of what the query asks for.",
    )
    add_band(
        slide,
        8.9,
        1.2,
        3.8,
        4.8,
        "Setup C: agentic",
        "Framework-inspired orchestration.\nSpecialist steps induce or retrieve components.\nE.text search uses extract-from-chunks plus reduce selection over report chunks.",
    )


def add_experiment_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Three Experiments", "Each experiment asks for a different missing framework component.")
    add_band(
        slide,
        0.7,
        1.05,
        12.0,
        1.65,
        "Experiment 1: D:(D.text,_)  P:(M,_,{(?,_)[3]})",
        "Given a report and a mathematical pattern definition M, retrieve three report excerpts that instantiate the pattern.",
    )
    add_band(
        slide,
        0.7,
        2.95,
        12.0,
        1.65,
        "Experiment 2: D:(D.text,_)  P:(_,_,{(?,E.tab)})",
        "Given a report and a compact table, recover the report excerpt that expresses the same evidence as the table.",
    )
    add_band(
        slide,
        0.7,
        4.85,
        12.0,
        1.65,
        "Experiment 3: D:(_,_)  P:(?,?,{(E1.text,E1.tab),(E2.text,E2.tab),...})",
        "Given paired text/table examples, induce the general mathematical rule M and natural-language description T.",
    )


def add_models_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Models Tested", "Gemma and Qwen thinking/mixed runs are intentionally excluded from this deck.")
    add_bullets(
        slide,
        0.85,
        1.25,
        11.8,
        4.7,
        [
            "gpt-4.1: strongest direct baseline and the judge model used for evaluation.",
            "gpt-oss-high: OpenRouter gpt-oss-120b with high reasoning effort.",
            "gpt-oss-medium: OpenRouter gpt-oss-120b with medium reasoning effort.",
            "qwen3-32B non-reasoning: OpenRouter Qwen with reasoning disabled and providers pinned.",
            "llama3-70B: Llama 3.3 70B comparison run.",
        ],
        24,
    )


def add_audit_overview_slide(prs: Presentation, rows: list[dict[str, Any]]) -> None:
    counts = defaultdict(int)
    gpt_rows = [row for row in rows if row["model"] == "gpt-4.1"]
    for row in rows:
        counts[int(row["score"])] += 1
    gpt_counts = defaultdict(int)
    for row in gpt_rows:
        gpt_counts[int(row["score"])] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Judge Scoring Audit", "Rationales were not stored; score slides use gpt-4.1 examples for apples-to-apples comparison.")
    body = [
        "Score 0 is assigned before an LLM judge call when the predicted excerpt cannot be recovered or snapped to the report.",
        "Scores 1-5 are judge alignment scores: 1 unrelated, 3 partial, 5 clear and complete.",
        f"Full reconstructed audit export: {OUT_AUDIT.name}",
        "Retained robustness score counts, all models: " + ", ".join(f"{score}: {counts[score]}" for score in range(6)),
        "Displayed gpt-4.1 score counts: " + ", ".join(f"{score}: {gpt_counts[score]}" for score in range(6)),
    ]
    add_bullets(slide, 0.75, 1.2, 11.9, 4.8, body, 23)


def add_conclusions_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusions", "The setup ranking depends on task shape, but the macro pattern is clear.")
    add_bullets(
        slide,
        0.75,
        1.1,
        11.9,
        5.4,
        [
            "Setup C is the strongest overall setup: it dominates E1 retrieval and usually wins the macro averages for capable models.",
            "Setup A remains a strong baseline, especially for direct E2 recovery and for weaker/non-reasoning models on natural-language induction.",
            "B-query is diagnostically useful: formal framework notation alone can be a burden when the model must infer the task semantics from rules.",
            "B-resolved is the fairer structured-prompt baseline because it keeps the framework but supplies the natural-language task resolution.",
            "At larger scale, agentic decomposition starts paying off; at smaller or less reliable reasoning scale, model specialization and instruction-following limits are exposed.",
        ],
        23,
    )


def build_deck() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    rows = audit_rows()
    write_audit_csv(rows)
    gpt_rows = [row for row in rows if row["model"] == "gpt-4.1"]
    by_score_counts = defaultdict(int)
    for row in gpt_rows:
        by_score_counts[int(row["score"])] += 1
    examples = choose_audit_examples(rows, "gpt-4.1")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_setup_slide(prs)
    add_experiment_slide(prs)
    add_models_slide(prs)
    make_result_slides(prs)
    add_audit_overview_slide(prs, rows)
    for score in range(5, -1, -1):
        add_audit_slide(prs, score, examples[score], by_score_counts[score])
    add_conclusions_slide(prs)
    prs.save(OUT_PPTX)


if __name__ == "__main__":
    build_deck()
    print(OUT_PPTX)
    print(OUT_AUDIT)
